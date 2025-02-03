from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Создаем презентацию
prs = Presentation()

# Стиль для заголовков
title_style = {
    'font_size': Pt(44),
    'color': RGBColor(59, 89, 152),
    'bold': True
}

# Стиль для основного текста
content_style = {
    'font_size': Pt(24),
    'color': RGBColor(0, 0, 0)
}

# ---------------------------
# Слайд 1: Титульный слайд
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[0])

title = slide.shapes.title
title.text = "Evomap\nСтратегия территориального развития"
title.text_frame.paragraphs[0].font.size = title_style['font_size']
title.text_frame.paragraphs[0].font.color.rgb = title_style['color']

subtitle = slide.placeholders[1]
subtitle.text = "Авторы: Гоша и Булат"
subtitle.text_frame.paragraphs[0].font.size = Pt(28)

# ---------------------------
# Слайд 2: Описание идеи
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Основная концепция"

content = [
    "• Стратегия с управлением городами и ресурсами",
    "• Процедурная генерация мира (алгоритм шума Перлина)",
    "• Динамическая экономика, зависящая от ландшафта",
    "• Система строительства железных дорог между городами"
]

text_frame = slide.shapes.placeholders[1].text_frame
for item in content:
    p = text_frame.add_paragraph()
    p.text = item
    p.font.size = content_style['font_size']

# ---------------------------
# Слайд 3: Реализация
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Техническая реализация"

content = [
    "Ключевые компоненты системы:",
    "- Класс Town: управление зданиями и населением",
    "- Класс Map: работа с чанками 32x32 пикселя",
    "- Класс Rail: логика соединения городов",
    "",
    "Особенности реализации:",
    "• Формулы расчета эффективности добычи ресурсов",
    "• Алгоритм автоматической отрисовки рельсов",
    "• Адаптивная система подсказок в интерфейсе"
]

text_frame = slide.shapes.placeholders[1].text_frame
for item in content:
    p = text_frame.add_paragraph()
    p.text = item
    p.font.size = content_style['font_size']

# ---------------------------
# Слайд 4: Технологии
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Используемые технологии"

content = [
    "Основной стек разработки:",
    "• Python 3.9+ - ядро системы",
    "• Pygame - визуализация и интерфейс",
    "• Numpy - обработка карты",
    "",
    "Запуск проекта:",
    "pip install pygame numpy\npython main.py"
]

text_frame = slide.shapes.placeholders[1].text_frame
for item in content:
    p = text_frame.add_paragraph()
    p.text = item
    p.font.size = content_style['font_size']

# ---------------------------
# Слайд 5: Игровой процесс
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Особенности игрового процесса"

content = [
    "Основные механики:",
    "1. Разведка и анализ территории",
    "2. Строительство и развитие городов",
    "3. Оптимизация логистических маршрутов",
    "4. Управление ресурсами и населением",
    "",
    "Ключевые показатели:",
    "- Коэффициенты добычи ресурсов",
    "- Уровни развития построек",
    "- Стабильность снабжения городов"
]

text_frame = slide.shapes.placeholders[1].text_frame
for item in content:
    p = text_frame.add_paragraph()
    p.text = item
    p.font.size = content_style['font_size']

# ---------------------------
# Слайд 6: Перспективы
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Планы развития"

content = [
    "Основные направления развития:",
    "• Введение торговой системы между городами",
    "• Добавление случайных событий и катастроф",
    "• Реализация сетевого мультиплеера",
    "• Расширение системы технологического прогресса",
    "",
    "Долгосрочные цели:",
    "- Динамическая смена климата",
    "- Система политических альянсов",
    "- Генерация уникальных событий"
]

text_frame = slide.shapes.placeholders[1].text_frame
for item in content:
    p = text_frame.add_paragraph()
    p.text = item
    p.font.size = content_style['font_size']

# Сохраняем презентацию
prs.save('Evomap_Text_Presentation.pptx')
print("Текстовая презентация успешно создана!")